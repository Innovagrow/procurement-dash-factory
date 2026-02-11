"""
Export Dashboards to PDF and PowerPoint

Generate professional reports for sharing
"""
from __future__ import annotations
from pathlib import Path
from typing import Any
import json
from datetime import datetime


def export_to_pdf(
    dashboard_html: str | Path,
    output_path: str | Path,
    title: str | None = None
) -> Path:
    """
    Export dashboard to PDF
    
    Requires: pip install playwright weasyprint
    """
    output_path = Path(output_path)
    
    try:
        from weasyprint import HTML, CSS
    except ImportError:
        raise ImportError("Install weasyprint: pip install weasyprint")
    
    # Read HTML
    if isinstance(dashboard_html, Path):
        html_content = dashboard_html.read_text(encoding='utf-8')
    else:
        html_content = dashboard_html
    
    # Add print styles
    css = CSS(string='''
        @page {
            size: A4 landscape;
            margin: 1cm;
        }
        body {
            font-family: Arial, sans-serif;
        }
        .no-print {
            display: none !important;
        }
    ''')
    
    # Generate PDF
    HTML(string=html_content).write_pdf(output_path, stylesheets=[css])
    
    return output_path


def export_to_powerpoint(
    plan: dict[str, Any],
    dataset_code: str,
    output_path: str | Path
) -> Path:
    """
    Export dashboard to PowerPoint presentation
    
    Requires: pip install python-pptx
    """
    output_path = Path(output_path)
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = plan.get('dataset', {}).get('title', dataset_code)
    subtitle.text = f"AI-Enhanced Analytics Report\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # Executive Summary slide
    if 'ai_insights' in plan and 'narrative' in plan['ai_insights']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Summary"
        
        tf = body_shape.text_frame
        tf.text = "Key Findings:"
        
        # Add insights as bullets
        insights = plan['ai_insights'].get('insights', [])
        for insight in insights[:5]:  # Top 5 insights
            p = tf.add_paragraph()
            p.text = f"{insight.get('title', '')}"
            p.level = 1
    
    # Key Metrics slide
    if 'ai_insights' in plan and 'key_metrics' in plan['ai_insights']:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Key Metrics"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Add metrics in a grid
        metrics = plan['ai_insights']['key_metrics']
        cols = 2
        rows = (len(metrics) + cols - 1) // cols
        
        left_margin = 0.5
        top_margin = 1.5
        box_width = 4.25
        box_height = 1.5
        spacing = 0.5
        
        for i, metric in enumerate(metrics):
            row = i // cols
            col = i % cols
            
            left = Inches(left_margin + col * (box_width + spacing))
            top = Inches(top_margin + row * (box_height + spacing))
            
            box = slide.shapes.add_textbox(left, top, Inches(box_width), Inches(box_height))
            frame = box.text_frame
            
            p = frame.paragraphs[0]
            p.text = metric.get('name', '')
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            
            p2 = frame.add_paragraph()
            value = metric.get('value', 0)
            p2.text = f"{value:,.0f}" if isinstance(value, (int, float)) else str(value)
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.alignment = PP_ALIGN.CENTER
    
    # Add page for each dashboard section
    for page in plan.get('pages', []):
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = page.get('title', 'Analysis')
        
        tf = body_shape.text_frame
        tf.text = page.get('description', '')
        
        # Add visualizations as bullets
        for visual in page.get('visuals', [])[:3]:  # Limit to 3 per slide
            p = tf.add_paragraph()
            p.text = visual.get('title', '')
            p.level = 1
            
            if visual.get('description'):
                p2 = tf.add_paragraph()
                p2.text = visual.get('description', '')
                p2.level = 2
    
    # Save presentation
    prs.save(str(output_path))
    return output_path


def export_to_excel(
    data: dict[str, Any],
    output_path: str | Path
) -> Path:
    """
    Export data to Excel workbook
    
    Requires: pip install openpyxl
    """
    import pandas as pd
    
    output_path = Path(output_path)
    
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        # Summary sheet
        summary_data = {
            'Metric': [],
            'Value': []
        }
        
        if 'ai_insights' in data and 'key_metrics' in data['ai_insights']:
            for metric in data['ai_insights']['key_metrics']:
                summary_data['Metric'].append(metric.get('name', ''))
                summary_data['Value'].append(metric.get('value', ''))
        
        pd.DataFrame(summary_data).to_excel(writer, sheet_name='Summary', index=False)
        
        # Insights sheet
        if 'ai_insights' in data and 'insights' in data['ai_insights']:
            insights_data = {
                'Title': [],
                'Description': [],
                'Type': [],
                'Severity': []
            }
            
            for insight in data['ai_insights']['insights']:
                insights_data['Title'].append(insight.get('title', ''))
                insights_data['Description'].append(insight.get('description', ''))
                insights_data['Type'].append(insight.get('type', ''))
                insights_data['Severity'].append(insight.get('severity', ''))
            
            pd.DataFrame(insights_data).to_excel(writer, sheet_name='Insights', index=False)
    
    return output_path
